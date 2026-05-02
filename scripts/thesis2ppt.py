#!/usr/bin/env python3
"""
Thesis to PPT Generator - 毕业论文答辩PPT生成器

Supports: PDF, DOCX, TXT, Markdown input formats.
Outputs: PPTX with academic blue theme styling.

Usage:
  # Parse only (extract structure as JSON)
  python thesis2ppt.py --parse-only thesis.pdf --output-json structure.json

  # Generate PPT
  python thesis2ppt.py thesis.pdf --output defense.pptx \\
      --title "论文标题" --author "张三" --advisor "李教授" --university "XX大学"
"""

import argparse
import json
import os
import re
import sys
import zipfile
import tempfile
import shutil
import subprocess
from pathlib import Path
from typing import Optional

# ---------------------------------------------------------------------------
# Structure extraction patterns (Chinese thesis focused)
# ---------------------------------------------------------------------------

CHAPTER_PATTERNS = [
    # 第X章 / 第X章 标题
    re.compile(r"^第[一二三四五六七八九十\d]+章\s*[\.、\s]?\s*(.*)"),
    # 一、二、三、... (top-level section)
    re.compile(r"^[一二三四五六七八九十]+[、.]\s*(.*)"),
    # (一)(二)(三)...
    re.compile(r"^[（(][一二三四五六七八九十]+[)）]\s*(.*)"),
    # 1. / 1.1 / 1.1.1
    re.compile(r"^(\d+(?:\.\d+)*)\s+([^\d].*)"),
    # Abstract / 摘要
    re.compile(r"^(摘\s*要|Abstract|ABSTRACT)\s*$"),
    # Keywords / 关键词
    re.compile(r"^(关键词|Keywords|KEYWORDS)\s*[：:]?"),
    # 参考文献 / References
    re.compile(r"^(参考文献|References|REFERENCES)\s*$"),
    # 致谢 / Acknowledgements
    re.compile(r"^(致\s*谢|Acknowledgements?|ACKNOWLEDGEMENTS?)\s*$"),
    # 目录 / Table of Contents
    re.compile(r"^(目\s*录|Table\s+of\s+Contents)\s*$"),
    # 附录 / Appendix
    re.compile(r"^(附\s*录|Appendix)\s*[：:]*\s*(.*)"),
    # 绪论 / 引言 / Introduction
    re.compile(r"^(绪\s*论|引\s*言|Introduction|INTRODUCTION)\s*$"),
    # 结论 / Conclusion
    re.compile(r"^(结\s*论|总结|Conclusions?|CONCLUSIONS?)\s*$"),
]

SECTION_INDICATORS = [
    re.compile(r"^(\d+(?:\.\d+)*)\s+(.*)"),         # 1.1 / 1.1.1 Title
    re.compile(r"^[（(]([一二三四五六七八九十]+)[)）]\s*(.*)"),  # (一) Title
    re.compile(r"^([一二三四五六七八九十]+)[、.]\s*(.*)"),      # 一、Title
]


def is_chapter_title(line: str) -> Optional[tuple[int, str]]:
    """Check if a line is a chapter/section title. Returns (level, title) or None."""
    line = line.strip()
    if not line or len(line) > 80:
        return None

    for i, pattern in enumerate(CHAPTER_PATTERNS):
        m = pattern.match(line)
        if m:
            if i == 0:  # 第X章
                title = m.group(1) if m.group(1) else line
                return (1, title)
            elif i == 1:  # 一、
                return (2, m.group(1) if m.group(1) else line)
            elif i == 2:  # (一)
                return (3, m.group(1) if m.group(1) else line)
            elif i == 3:  # 1. / 1.1 / 1.1.1
                level = m.group(1).count(".") + 1
                return (min(level, 4), m.group(2).strip())
            elif i == 4:  # 摘要
                return (1, line)
            elif i == 8:  # 目录
                return (0, line)
            elif i == 6:  # 参考文献
                return (1, "参考文献")
            elif i == 7:  # 致谢
                return (1, "致谢")
            elif i == 9:  # 附录
                rest = m.group(2) if m.lastindex and m.group(2) else ""
                return (1, f"附录 {rest}".strip())
            elif i == 10:  # 绪论/引言
                return (1, line)
            elif i == 11:  # 结论
                return (1, line)
            else:
                return (2, line)
    return None


def detect_chinese_number(text: str) -> int:
    """Convert Chinese number to integer (一 -> 1, 二十 -> 20, etc.)"""
    mapping = {
        "一": 1, "二": 2, "三": 3, "四": 4, "五": 5,
        "六": 6, "七": 7, "八": 8, "九": 9, "十": 10,
    }
    if text in mapping:
        return mapping[text]
    if text.startswith("十"):
        rest = text[1:]
        return 10 + (mapping.get(rest, 0) if rest else 0)
    if text.endswith("十"):
        return mapping.get(text[0], 0) * 10
    return 0


# ---------------------------------------------------------------------------
# Parsers
# ---------------------------------------------------------------------------

class ThesisParser:
    """Parses thesis files into structured content."""

    def parse(self, filepath: str) -> list[dict]:
        """Parse a thesis file and return structured sections."""
        ext = Path(filepath).suffix.lower()
        if ext == ".pdf":
            return self._parse_pdf(filepath)
        elif ext == ".docx":
            return self._parse_docx(filepath)
        elif ext in (".txt", ".md", ".markdown"):
            return self._parse_text(filepath)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def _parse_pdf(self, filepath: str) -> list[dict]:
        """Parse PDF thesis."""
        try:
            import pdfplumber
        except ImportError:
            raise ImportError("pdfplumber is required. Install: pip install pdfplumber")

        all_text = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    all_text.append(text)

        full_text = "\n".join(all_text)
        return self._extract_structure(full_text)

    def _parse_docx(self, filepath: str) -> list[dict]:
        """Parse Word thesis."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required. Install: pip install python-docx")

        doc = Document(filepath)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Use Word styles to infer heading level
            style_name = para.style.name if para.style else ""
            if style_name.startswith("Heading") or style_name.startswith("heading"):
                try:
                    level = int(style_name.replace("Heading", "").replace("heading", "").strip())
                except ValueError:
                    level = 1
                paragraphs.append({"text": text, "is_heading": True, "level": min(level, 4)})
            else:
                paragraphs.append({"text": text, "is_heading": False, "level": 0})

        return self._structure_from_paragraphs(paragraphs, filepath)

    def _parse_text(self, filepath: str) -> list[dict]:
        """Parse plain text / Markdown thesis."""
        with open(filepath, "r", encoding="utf-8") as f:
            content = f.read()

        if filepath.endswith(".md"):
            return self._parse_markdown(content)
        return self._extract_structure(content)

    def _parse_markdown(self, content: str) -> list[dict]:
        """Parse Markdown with heading structure."""
        lines = content.split("\n")
        sections = []
        current_section = None

        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            heading_match = re.match(r"^(#{1,6})\s+(.*)", stripped)
            if heading_match:
                if current_section:
                    sections.append(current_section)
                level = len(heading_match.group(1))
                title = heading_match.group(2).strip()
                current_section = {
                    "title": title,
                    "level": min(level, 4),
                    "content": "",
                    "subsections": [],
                }
            elif current_section:
                current_section["content"] += stripped + "\n"

        if current_section:
            sections.append(current_section)
        return sections

    def _extract_structure(self, text: str) -> list[dict]:
        """Extract section structure from raw text."""
        lines = text.split("\n")
        sections = []
        current_section = None
        current_content = []

        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue

            title_info = is_chapter_title(stripped)
            if title_info:
                # Save previous section
                if current_section is not None:
                    current_section["content"] = "\n".join(current_content)
                    sections.append(current_section)

                level, title = title_info
                current_section = {
                    "title": title,
                    "level": level,
                    "content": "",
                    "subsections": [],
                }
                current_content = []
            elif current_section is not None:
                current_content.append(stripped)

        if current_section is not None:
            current_section["content"] = "\n".join(current_content)
            sections.append(current_section)

        return sections

    def _structure_from_paragraphs(self, paragraphs: list[dict], filepath: str) -> list[dict]:
        """Build structure from docx paragraphs with style info."""
        sections = []
        current_section = None
        current_content = []

        for para in paragraphs:
            if para["is_heading"]:
                if current_section is not None:
                    current_section["content"] = "\n".join(current_content)
                    sections.append(current_section)

                current_section = {
                    "title": para["text"],
                    "level": para["level"],
                    "content": "",
                    "subsections": [],
                }
                current_content = []
            elif current_section is not None:
                current_content.append(para["text"])
            else:
                # Text before any heading — preamble
                if not sections and not current_section:
                    current_section = {
                        "title": "论文信息",
                        "level": 0,
                        "content": "",
                        "subsections": [],
                    }
                    current_content.append(para["text"])

        if current_section is not None:
            current_section["content"] = "\n".join(current_content)
            sections.append(current_section)

        return sections


# ---------------------------------------------------------------------------
# PPT Builder
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Image extraction from DOCX
# ---------------------------------------------------------------------------

def convert_emf_to_png(emf_path: str, png_path: str) -> bool:
    """Convert an EMF file to PNG using Windows GDI via PowerShell.

    Returns True on success, False on failure.
    """
    try:
        ps_script = (
            f"Add-Type -AssemblyName System.Drawing;"
            f"$img = [System.Drawing.Imaging.Metafile]::FromFile('{emf_path}');"
            f"$img.Save('{png_path}', [System.Drawing.Imaging.ImageFormat]::Png);"
            f"$img.Dispose()"
        )
        result = subprocess.run(
            ["powershell.exe", "-NoProfile", "-Command", ps_script],
            capture_output=True, text=True, timeout=15,
        )
        return os.path.exists(png_path) and os.path.getsize(png_path) > 0
    except Exception:
        return False


def ensure_image_png(image_dir: str, filename: str) -> Optional[str]:
    """Ensure an image is in PNG format suitable for pptx.

    Returns the PNG file path, or None if conversion is not possible.
    """
    src_path = os.path.join(image_dir, filename)
    if not os.path.exists(src_path):
        return None

    ext = os.path.splitext(filename)[1].lower()
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif"):
        return src_path

    if ext == ".emf":
        png_name = os.path.splitext(filename)[0] + ".png"
        png_path = os.path.join(image_dir, png_name)
        if os.path.exists(png_path):
            return png_path
        if convert_emf_to_png(src_path, png_path):
            return png_path

    return None

def extract_images_from_docx(docx_path: str, output_dir: str) -> list[str]:
    """Extract embedded images from a DOCX file.

    Returns list of extracted image file paths.
    """
    os.makedirs(output_dir, exist_ok=True)
    extracted = []

    with zipfile.ZipFile(docx_path, "r") as z:
        # Find all image files in word/media/ (skip directory entries)
        image_names = [
            n for n in z.namelist()
            if n.startswith("word/media/") and not n.endswith("/")
        ]

        for name in image_names:
            basename = os.path.basename(name)
            out_path = os.path.join(output_dir, basename)
            with z.open(name) as src, open(out_path, "wb") as dst:
                dst.write(src.read())
            extracted.append(out_path)

    return extracted


def find_image_references(docx_path: str) -> list[dict]:
    """Find image references in DOCX document.xml with surrounding context.

    Returns list of dicts with:
        filename, figure_caption (图注 like "图3.4 ..."),
        context (surrounding descriptive text), prev_text, next_text.
    """
    refs = []
    with zipfile.ZipFile(docx_path, "r") as z:
        # Read relationships to map rId -> filename
        rels_xml = z.read("word/_rels/document.xml.rels").decode("utf-8")
        rid_to_file = {}
        for m in re.finditer(
            r'Id="(rId\d+)"[^>]*Target="([^"]*)"', rels_xml
        ):
            rid_to_file[m.group(1)] = m.group(2)

        # Read document XML and split into paragraphs
        doc_xml = z.read("word/document.xml").decode("utf-8")
        paragraphs = re.split(r"</w:p>", doc_xml)

        for i, para in enumerate(paragraphs):
            blips = re.findall(r'r:embed="(rId\d+)"', para)
            if not blips:
                continue

            # Extract text from this paragraph and neighbors
            def _extract_text(p):
                return "".join(re.findall(r"<w:t[^>]*>([^<]*)</w:t>", p)).strip()

            para_text = _extract_text(para)
            prev_text = _extract_text(paragraphs[i - 1]) if i > 0 else ""
            next_text = _extract_text(paragraphs[i + 1]) if i + 1 < len(paragraphs) else ""

            # Identify the figure caption: text matching "图X.X ..." pattern
            figure_caption = ""
            # Check next paragraph first (most common location for captions)
            for candidate in [next_text, prev_text, para_text]:
                if re.match(r"图\s*\d+", candidate):
                    figure_caption = candidate
                    break
            # Fallback: use anything that looks like a caption
            if not figure_caption:
                for candidate in [next_text, para_text, prev_text]:
                    if candidate and len(candidate) < 120 and (
                        "图" in candidate or "Figure" in candidate
                    ):
                        figure_caption = candidate
                        break
            # Last resort
            if not figure_caption:
                figure_caption = para_text or next_text or prev_text

            # Context = descriptive text around the image (not the caption itself)
            context = para_text if para_text != figure_caption else prev_text

            for rid in blips:
                if rid in rid_to_file:
                    refs.append({
                        "embed_id": rid,
                        "filename": rid_to_file[rid],
                        "caption": figure_caption,   # The actual 图注
                        "context": context,          # Surrounding description
                        "prev_text": prev_text,
                        "next_text": next_text,
                    })

    return refs


class PPTBuilder:
    """Generates PPTX files with academic blue styling."""

    def __init__(self, config_path: Optional[str] = None):
        self.config = self._load_config(config_path)
        self._setup_theme()

    def _load_config(self, config_path: Optional[str]) -> dict:
        """Load style configuration."""
        if config_path and os.path.exists(config_path):
            with open(config_path, "r", encoding="utf-8") as f:
                return json.load(f)

        # Default config
        script_dir = Path(__file__).parent.parent
        default_config = script_dir / "templates" / "academic_blue.json"
        if default_config.exists():
            with open(default_config, "r", encoding="utf-8") as f:
                return json.load(f)

        return {
            "colors": {
                "primary": "003366",
                "accent": "0066CC",
                "text": "333333",
                "background": "FFFFFF",
                "light_bg": "E8F0FE",
                "white": "FFFFFF",
            },
            "fonts": {
                "title": "SimHei",
                "body": "SimSun",
                "size_title": 32,
                "size_subtitle": 24,
                "size_section": 28,
                "size_body": 18,
                "size_small": 14,
                "size_footer": 10,
            },
            "slide_width_cm": 25.4,
            "slide_height_cm": 19.05,
        }

    def _setup_theme(self):
        """Initialize theme attributes from config."""
        self.c = self.config["colors"]
        self.f = self.config["fonts"]
        w_cm = self.config.get("slide_width_cm", 25.4)
        h_cm = self.config.get("slide_height_cm", 19.05)
        self.w_cm = w_cm  # Keep cm for slide dimensions
        self.h_cm = h_cm
        self.w = w_cm / 2.54  # Convert to inches for shape positioning
        self.h = h_cm / 2.54

    def _hex_to_rgb(self, hex_color: str) -> tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    def _add_textbox(self, slide, left, top, width, height, text, font_name,
                     font_size, color_hex, bold=False, alignment=None):
        """Add a text box to a slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*self._hex_to_rgb(color_hex))
        if alignment is not None:
            p.alignment = alignment
        return tf

    def create_cover_slide(self, prs, title: str, author: str = "",
                           advisor: str = "", university: str = "",
                           date_str: str = ""):
        """Create the cover/title slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Top decorative bar
        left, top, width, height = 0, 0, self.w, 1.2
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height)  # Rectangle
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["primary"]))
        shape.line.fill.background()

        # University name
        if university:
            self._add_textbox(
                slide, 1.5, 1.6, self.w - 3, 0.6,
                university, self.f["body"], self.f["size_small"],
                self.c["text"], bold=False,
                alignment=PP_ALIGN.CENTER,
            )

        # Title
        self._add_textbox(
            slide, 2.0, 2.5, self.w - 4, 1.8,
            title, self.f["title"], self.f["size_title"],
            self.c["primary"], bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # Decorative line below title
        shape = slide.shapes.add_shape(
            1, Inches(3.5), Inches(4.5), Inches(self.w - 7), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["accent"]))
        shape.line.fill.background()

        # Info section
        y_start = 5.5
        if author:
            self._add_textbox(
                slide, 4.0, y_start, self.w - 8, 0.5,
                f"答辩人：{author}", self.f["body"], self.f["size_body"],
                self.c["text"], alignment=PP_ALIGN.CENTER,
            )
            y_start += 0.6
        if advisor:
            self._add_textbox(
                slide, 4.0, y_start, self.w - 8, 0.5,
                f"指导教师：{advisor}", self.f["body"], self.f["size_body"],
                self.c["text"], alignment=PP_ALIGN.CENTER,
            )
            y_start += 0.6
        if date_str:
            self._add_textbox(
                slide, 4.0, y_start + 0.3, self.w - 8, 0.5,
                date_str, self.f["body"], self.f["size_body"],
                self.c["text"], alignment=PP_ALIGN.CENTER,
            )

        # Bottom decorative bar
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(self.h - 0.5), Inches(self.w), Inches(0.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["primary"]))
        shape.line.fill.background()

    def create_toc_slide(self, prs, items: list[str]):
        """Create table of contents slide."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Header bar
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(self.w), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["primary"]))
        shape.line.fill.background()

        self._add_textbox(
            slide, 0.8, 0.2, self.w - 1.6, 0.8,
            "目  录", self.f["title"], self.f["size_section"],
            self.c["white"], bold=True, alignment=PP_ALIGN.LEFT,
        )

        # TOC items with numbering
        y = 1.6
        for i, item in enumerate(items, 1):
            # Number circle
            circle = slide.shapes.add_shape(
                9, Inches(1.2), Inches(y + 0.05), Inches(0.35), Inches(0.35)  # Oval
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["primary"]))
            circle.line.fill.background()
            tf = circle.text_frame
            tf.paragraphs[0].text = str(i)
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(*self._hex_to_rgb(self.c["white"]))
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            self._add_textbox(
                slide, 1.8, y, self.w - 3.5, 0.45,
                item, self.f["body"], self.f["size_body"],
                self.c["text"], alignment=PP_ALIGN.LEFT,
            )
            y += 0.65

    def create_content_slide(self, prs, title: str, bullets: list[str],
                              section_num: str = "", image_path: str = None,
                              image_caption: str = ""):
        """Create a content slide with left-image + right-text wrap-around layout.

        Args:
            image_path: Optional path to an image file to embed on the left side.
            image_caption: Short description shown below the image.
        """
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Top title bar
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(self.w), Inches(1.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["primary"]))
        shape.line.fill.background()

        # Title text
        display_title = f"{section_num}  {title}" if section_num else title
        self._add_textbox(
            slide, 0.5, 0.15, self.w - 1, 0.7,
            display_title, self.f["title"], self.f["size_subtitle"],
            self.c["white"], bold=True, alignment=PP_ALIGN.LEFT,
        )

        has_image = image_path and os.path.exists(image_path)

        if has_image:
            # ---- Left side: Image ----
            img_left = 0.6
            img_top = 1.4
            img_w = 4.2
            img_h = self.h - 3.0

            try:
                pic = slide.shapes.add_picture(
                    image_path,
                    Inches(img_left), Inches(img_top),
                    Inches(img_w), Inches(img_h),
                )
                # Scale to fit the reserved area, maintaining aspect ratio
                target_w_emu = int(img_w * 914400)
                target_h_emu = int(img_h * 914400)
                ratio = min(target_w_emu / pic.width, target_h_emu / pic.height, 1.0)
                pic.width = int(pic.width * ratio)
                pic.height = int(pic.height * ratio)
                # Center image in its region
                if pic.width < target_w_emu:
                    offset_x = (target_w_emu - pic.width) // 2
                    pic.left = int(Inches(img_left)) + offset_x
                if pic.height < target_h_emu:
                    offset_y = (target_h_emu - pic.height) // 2
                    pic.top = int(Inches(img_top)) + offset_y
            except Exception:
                has_image = False

            # Image caption below
            if image_caption:
                short_cap = image_caption[:50] + ("..." if len(image_caption) > 50 else "")
                self._add_textbox(
                    slide, img_left, self.h - 1.3, img_w, 0.6,
                    short_cap, self.f["body"], self.f["size_small"],
                    self.c["accent"], alignment=PP_ALIGN.CENTER,
                )

            # ---- Right side: Text bullets ----
            text_left = 5.3
            text_width = self.w - 5.9
        else:
            # Full-width text, no image
            text_left = 1.0
            text_width = self.w - 2.0

        # Accent line separating image and text (only if image exists)
        if has_image:
            shape = slide.shapes.add_shape(
                1, Inches(5.1), Inches(1.4), Inches(0.03), Inches(self.h - 3.2)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["accent"]))
            shape.line.fill.background()

        # Bullet points on the right side
        y = 1.6
        for bullet in bullets:
            if not bullet.strip():
                continue
            # Bullet dot
            dot = slide.shapes.add_shape(
                9, Inches(text_left + 0.05), Inches(y + 0.06),
                Inches(0.12), Inches(0.12),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["accent"]))
            dot.line.fill.background()

            self._add_textbox(
                slide, text_left + 0.3, y, text_width - 0.4, 0.5,
                bullet.strip(), self.f["body"], self.f["size_body"],
                self.c["text"], alignment=PP_ALIGN.LEFT,
            )
            y += 0.6

        # Footer
        self._add_textbox(
            slide, 0.5, self.h - 0.35, self.w - 1, 0.3,
            "", self.f["body"], 8, self.c["text"], alignment=PP_ALIGN.CENTER,
        )

    def create_section_divider(self, prs, section_title: str, section_num: str = ""):
        """Create a section divider slide."""
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Full background with light color
        bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(self.w), Inches(self.h)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["light_bg"]))
        bg.line.fill.background()

        # Left accent bar
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(0.15), Inches(self.h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["primary"]))
        shape.line.fill.background()

        # Section number
        if section_num:
            self._add_textbox(
                slide, 2.0, 3.0, self.w - 4, 1.5,
                section_num, self.f["title"], 56,
                self.c["accent"], bold=True, alignment=PP_ALIGN.LEFT,
            )

        # Section title
        self._add_textbox(
            slide, 2.0, 5.0, self.w - 4, 1.2,
            section_title, self.f["title"], self.f["size_title"],
            self.c["primary"], bold=True, alignment=PP_ALIGN.LEFT,
        )

    def create_image_slide(self, prs, image_path: str, caption: str = "",
                           section_num: str = ""):
        """Create a slide with an embedded image and caption."""
        from pptx.util import Inches, Cm, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Top title bar
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(self.w), Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["primary"]))
        shape.line.fill.background()

        # Simplified caption as title
        short_caption = caption[:60] + ("..." if len(caption) > 60 else "")
        display_title = f"{section_num}  {short_caption}" if section_num else short_caption or "图表"
        self._add_textbox(
            slide, 0.5, 0.1, self.w - 1, 0.7,
            display_title, self.f["title"], self.f["size_subtitle"],
            self.c["white"], bold=True, alignment=PP_ALIGN.LEFT,
        )

        # Try to add the image
        try:
            if os.path.exists(image_path):
                # Calculate image placement (centered, max 80% of slide)
                max_img_w = self.w * 0.85
                max_img_h = self.h - 2.5  # Leave room for title bar + caption

                # Add picture with auto-scaling
                pic = slide.shapes.add_picture(
                    image_path,
                    Inches((self.w - max_img_w) / 2),
                    Inches(1.3),
                    Inches(max_img_w),
                    Inches(max_img_h),
                )
                # Maintain aspect ratio
                pic.width = int(pic.width)
                pic.height = int(pic.height)
                if pic.width / self.w > pic.height / max_img_h:
                    scale = (max_img_w * 914400) / pic.width  # 914400 EMU per inch
                else:
                    scale = (max_img_h * 914400) / pic.height
        except Exception:
            pass  # Skip if image can't be loaded

        # Footer caption
        if caption:
            self._add_textbox(
                slide, 1.0, self.h - 1.2, self.w - 2, 0.8,
                caption[:120], self.f["body"], self.f["size_small"],
                self.c["text"], alignment=PP_ALIGN.CENTER,
            )

    def create_ending_slide(self, prs, text: str = "谢谢！请各位老师批评指正"):
        """Create ending / Q&A slide."""
        from pptx.util import Inches
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(self.w), Inches(self.h)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(self.c["primary"]))
        bg.line.fill.background()

        self._add_textbox(
            slide, 3.0, self.h / 2 - 1.2, self.w - 6, 1.5,
            text, self.f["title"], 40,
            self.c["white"], bold=True, alignment=PP_ALIGN.CENTER,
        )

        # Subtitle
        self._add_textbox(
            slide, 3.0, self.h / 2 + 0.8, self.w - 6, 1.0,
            "Q & A", self.f["title"], 28,
            self.c["white"], bold=False, alignment=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# Content summarizer (simple extractive, AI enhancement can be layered on)
# ---------------------------------------------------------------------------

def summarize_section(section: dict, max_bullets: int = 5, max_chars: int = 55) -> list[str]:
    """Extract key bullet points from section content, optimized for PPT display.

    Each bullet is limited to ~max_chars characters to prevent text overflow.
    """
    content = section.get("content", "")
    if not content:
        return [section.get("title", "")]

    # Split into sentences
    raw = re.split(r"[。；;.\n]", content)
    sentences = [s.strip() for s in raw if len(s.strip()) > 8]

    if not sentences:
        return [section.get("title", "")]

    # Truncate each sentence to fit PPT text box
    trimmed = []
    for s in sentences:
        if len(s) > max_chars:
            # Try to break at comma/pause for cleaner truncation
            half = s[:max_chars]
            last_break = max(half.rfind("，"), half.rfind(","), half.rfind("、"), half.rfind(" "))
            if last_break > max_chars // 2:
                s = s[:last_break]
            else:
                s = s[:max_chars - 3]
            # Ensure sentence ends cleanly
            if not any(s.endswith(c) for c in ("，", ",", "、", "；", ":", "：", "的", "和", "与", "了")):
                pass  # OK as-is
        trimmed.append(s)

    if len(trimmed) <= max_bullets:
        return trimmed[:max_bullets]

    # Score and pick the most significant sentences
    key_indicators = [
        "研究", "方法", "结果", "分析", "提出", "设计", "实现",
        "实验", "数据", "模型", "算法", "系统", "技术", "理论",
        "问题", "解决", "改进", "优化", "创新", "贡献", "结论",
        "本文", "我们", "主要", "关键", "重要", "显著",
        "仿真", "测试", "验证", "协议", "信号", "电路", "调制",
    ]
    scored = []
    for s in trimmed:
        score = sum(1 for kw in key_indicators if kw in s)
        if 10 <= len(s) <= max_chars:
            score += 2
        elif len(s) > max_chars:
            score -= 1  # Penalize overly long sentences
        scored.append((score, s))

    scored.sort(key=lambda x: x[0], reverse=True)
    bullets = [s for _, s in scored[:max_bullets]]
    return bullets if bullets else trimmed[:max_bullets]


def map_sections_to_slides(sections: list[dict]) -> list[dict]:
    """Map extracted sections to PPT slide structure.

    Includes level-2 AND level-3 sections (二级+三级标题).
    If a level-2 section has no direct content, merges content from its
    immediate level-3 subsections.
    Skips appendix, references, acknowledgements, abstract, and TOC.
    """
    slides = []

    skip_keywords = [
        "附录", "appendix", "abstract", "目录", "参考文献", "references",
        "致谢", "acknowledgement", "acknowledgements",
    ]

    def _is_skip(title):
        return any(kw in title.lower() for kw in skip_keywords)

    # Build a parent-child map: find level-3 sections belonging to each level-2
    l2_to_l3 = {}  # l2_title -> [l3_sections]
    all_l3_titles = set()
    for i, sec in enumerate(sections):
        if sec.get("level") == 3:
            all_l3_titles.add(sec.get("title", "").strip())
            # Find preceding level-2 section
            for j in range(i - 1, -1, -1):
                if sections[j].get("level") == 2:
                    l2_title = sections[j].get("title", "").strip()
                    if l2_title not in l2_to_l3:
                        l2_to_l3[l2_title] = []
                    l2_to_l3[l2_title].append(sec)
                    break

    slide_count = 0
    for sec in sections:
        title = sec.get("title", "").strip()
        level = sec.get("level", 1)
        content = sec.get("content", "").strip()

        if _is_skip(title) or level not in (2, 3):
            continue

        # For level-2 with no content, merge from subsections
        if level == 2 and not content and title in l2_to_l3:
            merged_parts = []
            for l3 in l2_to_l3[title]:
                l3_content = l3.get("content", "").strip()
                if l3_content:
                    # Take first 1-2 sentences from each subsection
                    sentences = re.split(r"[。；;.]", l3_content)
                    merged_parts.extend(s[:120] for s in sentences[:2] if s.strip())
            content = "；".join(merged_parts[:6]) if merged_parts else ""

        if not content:
            continue

        slide_count += 1
        section_num = f"{slide_count:02d}"
        slides.append({
            "type": "content",
            "title": title,
            "section_num": section_num,
            "content": content,
            "level": level,
        })

    return slides


# ---------------------------------------------------------------------------
# Main entry points
# ---------------------------------------------------------------------------

def parse_thesis(filepath: str) -> dict:
    """Parse thesis and return structured data as JSON-serializable dict."""
    parser = ThesisParser()
    sections = parser.parse(filepath)
    return {
        "file": filepath,
        "section_count": len(sections),
        "sections": sections,
    }


def generate_ppt(
    filepath: str,
    output_path: str,
    config_path: Optional[str] = None,
    title: str = "",
    author: str = "",
    advisor: str = "",
    university: str = "",
    date_str: str = "",
    structure_json: Optional[str] = None,
    image_mapping_json: Optional[str] = None,
):
    """Generate PPTX from thesis file.

    Args:
        image_mapping_json: Optional JSON file with Claude-reviewed image mappings.
            Format: {"mappings": [{"section_title": "...", "image": "file.png"}, ...]}
            When provided, uses these mappings instead of auto-matching.
    """
    from pptx import Presentation

    # Parse thesis
    parser = ThesisParser()
    sections = parser.parse(filepath)

    # If custom structure JSON provided, use it; otherwise auto-map
    if structure_json:
        with open(structure_json, "r", encoding="utf-8") as f:
            slide_plan = json.load(f)
    else:
        slide_plan = map_sections_to_slides(sections)

    # Build PPT
    builder = PPTBuilder(config_path)
    prs = Presentation()
    prs.slide_width = int(builder.w_cm * 360000)  # EMU: 1cm = 360000 EMU
    prs.slide_height = int(builder.h_cm * 360000)

    # Determine thesis title
    if not title:
        for sec in sections:
            if sec.get("title") and len(sec["title"]) > 3:
                title = sec["title"]
                break
        if not title:
            title = "毕业论文答辩"

    # Cover slide
    builder.create_cover_slide(prs, title, author, advisor, university, date_str)

    # TOC — use level-1 section titles (chapter titles like "第一章 绪论")
    toc_skip = ["附录", "appendix", "abstract", "目录", "参考文献", "references",
                "致谢", "acknowledgement"]
    toc_items = []
    for sec in sections:
        title = sec.get("title", "").strip()
        level = sec.get("level", 0)
        if level != 1:
            continue
        if any(kw in title.lower() for kw in toc_skip):
            continue
        # Clean up the title: remove leading number+whitespace, keep chapter info
        cleaned = re.sub(r"^\d+\s+", "", title)  # Remove standalone number prefix
        if cleaned and len(cleaned) > 1:
            toc_items.append(cleaned)

    if toc_items:
        builder.create_toc_slide(prs, toc_items)

    # Content slides data (already filtered to level-2 in map_sections_to_slides)
    content_slides_data = list(slide_plan)

    # Extract images from DOCX if available
    image_dir = None
    image_refs = []
    if filepath.lower().endswith(".docx"):
        try:
            # Use temp dir to avoid Chinese path encoding issues
            image_dir = os.path.join(tempfile.gettempdir(), "thesis_ppt_images")
            extract_images_from_docx(filepath, image_dir)
            image_refs = find_image_references(filepath)
            # Convert EMF images to PNG for pptx compatibility
            for ref in image_refs:
                fname = ref.get("filename", "")
                if fname.lower().endswith(".emf"):
                    ensure_image_png(image_dir, os.path.basename(fname))
        except Exception as e:
            print(f"  [info] Image extraction skipped: {e}")

    # Load image mapping (Claude-reviewed or user-matched)
    reviewed_mapping = {}  # section_title -> {"image": str, "images": [str]}
    if image_mapping_json and os.path.exists(image_mapping_json):
        with open(image_mapping_json, "r", encoding="utf-8") as f:
            mapping_data = json.load(f)
        for m in mapping_data.get("mappings", []):
            reviewed_mapping[m["section_title"]] = m

    # Build filename→ref lookup for quick access
    filename_to_ref = {}
    for ref in image_refs:
        fname = os.path.basename(ref.get("filename", ""))
        if fname:
            filename_to_ref[fname] = ref

    used_images = set()

    # Content slides: each section gets one paired image (left) + text (right)
    # Extra images for a section become separate image-only slides after the content slide
    for item in content_slides_data:
        sec_title = item.get("title", "")
        bullets = summarize_section(item, max_bullets=3)

        # Get all images for this section
        mapped_imgs = []
        if reviewed_mapping:
            rm = reviewed_mapping.get(sec_title, {})
            first_img = rm.get("image")
            extra_imgs = rm.get("images") or []
            if first_img:
                mapped_imgs = [first_img] + [im for im in extra_imgs if im != first_img]
        elif image_refs:
            # Fallback: auto-match single image
            def _best_image_for_section(title: str, refs: list[dict]):
                best, best_score = None, -1
                title_words = re.findall(r"[一-鿿]{2,}", title)
                for r in refs:
                    combined = r.get("caption", "") + r.get("prev_text", "") + r.get("next_text", "")
                    score = sum(len(w) for w in title_words if w in combined)
                    if re.search(r"图\s*\d+", combined):
                        score += 1
                    if score > best_score:
                        best_score, best = score, r
                return best if best_score > 0 else None

            best_ref = _best_image_for_section(sec_title, image_refs)
            if best_ref:
                fname = os.path.basename(best_ref.get("filename", ""))
                if fname not in used_images:
                    mapped_imgs = [fname]

        # First image goes into the content slide (left-image-right-text layout)
        image_path = None
        image_caption = ""
        if mapped_imgs and image_dir:
            first = mapped_imgs[0]
            if first in filename_to_ref:
                png_path = ensure_image_png(image_dir, first)
                if png_path and os.path.exists(png_path):
                    image_path = png_path
                    image_caption = filename_to_ref[first].get("caption", "")
                    used_images.add(first)

        builder.create_content_slide(
            prs, item.get("title", ""), bullets,
            item.get("section_num", ""),
            image_path=image_path, image_caption=image_caption,
        )

        # Extra images → standalone image slides
        for extra in mapped_imgs[1:]:
            if extra in used_images or extra not in filename_to_ref:
                continue
            png_path = ensure_image_png(image_dir, extra) if image_dir else None
            if png_path and os.path.exists(png_path):
                used_images.add(extra)
                ref = filename_to_ref[extra]
                builder.create_image_slide(
                    prs, png_path,
                    caption=ref.get("caption", ""),
                    section_num=item.get("section_num", ""),
                )

    # Acknowledgements slide
    builder.create_content_slide(
        prs, "致谢",
        ["感谢指导老师在论文完成过程中给予的悉心指导和帮助",
         "感谢各位评审老师在百忙之中抽出时间审阅本文",
         "感谢实验室的同学们在研究过程中的支持和协助",
         "感谢家人和朋友一直以来的关心和鼓励"],
        section_num="",
    )

    # Q&A ending slide
    builder.create_ending_slide(prs)

    # Save
    prs.save(output_path)
    return output_path


def main():
    # Fix Windows console encoding for Chinese output
    if sys.stdout.encoding != "utf-8":
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    parser = argparse.ArgumentParser(
        description="Thesis to Defense PPT Generator (毕业论文答辩PPT生成器)"
    )
    parser.add_argument("input", nargs="?", help="Path to thesis file")
    parser.add_argument("--parse-only", action="store_true",
                        help="Only parse and output structure JSON")
    parser.add_argument("--output-json", help="Output JSON path for parsed structure")
    parser.add_argument("--output", "-o", default="defense.pptx",
                        help="Output PPTX path (default: defense.pptx)")
    parser.add_argument("--config", help="Path to style config JSON")
    parser.add_argument("--title", help="Thesis title")
    parser.add_argument("--author", help="Author name")
    parser.add_argument("--advisor", help="Advisor name")
    parser.add_argument("--university", help="University name")
    parser.add_argument("--date", help="Date string")
    parser.add_argument("--structure", help="Custom structure JSON file")
    parser.add_argument("--image-mapping", help="Claude-reviewed image-to-section mapping JSON")

    args = parser.parse_args()

    if args.parse_only:
        if not args.input:
            print("Error: --parse-only requires an input file", file=sys.stderr)
            sys.exit(1)
        result = parse_thesis(args.input)
        if args.output_json:
            with open(args.output_json, "w", encoding="utf-8") as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            print(f"Structure saved to: {args.output_json}")
        else:
            json_str = json.dumps(result, ensure_ascii=False, indent=2)
            sys.stdout.reconfigure(encoding="utf-8")
            print(json_str)
        return

    if not args.input:
        parser.print_help()
        sys.exit(1)

    output = generate_ppt(
        filepath=args.input,
        output_path=args.output,
        config_path=args.config,
        title=args.title or "",
        author=args.author or "",
        advisor=args.advisor or "",
        university=args.university or "",
        date_str=args.date or "",
        structure_json=args.structure,
        image_mapping_json=args.image_mapping,
    )
    print(f"PPT generated: {output}")


if __name__ == "__main__":
    main()
